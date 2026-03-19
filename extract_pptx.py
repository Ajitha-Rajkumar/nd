"""
Extract text content from NPD.pptx and generate npd_data.js
Run: pip install python-pptx   (one-time)
Then: python extract_pptx.py
"""

import json
import re
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def extract_text_from_shape(shape):
    """Extract text from a shape, including tables and groups."""
    texts = []
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if text:
                texts.append(text)
    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        table = shape.table
        for row in table.rows:
            row_texts = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                if cell_text:
                    row_texts.append(cell_text)
            if row_texts:
                texts.append(" | ".join(row_texts))
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            texts.extend(extract_text_from_shape(s))
    return texts


def is_page_number(text):
    """Check if text is just a page/slide number."""
    return re.match(r'^\d{1,3}$', text.strip()) is not None


def clean_section_title(text):
    """Clean section headers like '2.\\tScope & Planning (Cont'd)' """
    # Remove leading numbering like "2.\t" or "2. " or "3.\t"
    cleaned = re.sub(r'^\d+\.\s*', '', text.strip())
    # Remove "(Cont'd)" / "(Cont'd)" / "(Continued)" with any apostrophe variant
    cleaned = re.sub(r"\s*\(Cont.?d\)", '', cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r"\s*\(Continued\)", '', cleaned, flags=re.IGNORECASE)
    # Remove extra whitespace
    cleaned = re.sub(r'\s+', ' ', cleaned)
    return cleaned.strip()


def detect_title(content_list):
    """Detect the title from slide content when no title placeholder exists."""
    if not content_list:
        return "", content_list

    # Strategy: first non-number, non-table-header item is usually the title
    for i, item in enumerate(content_list):
        text = item.strip()

        # Skip standalone numbers (page numbers)
        if is_page_number(text):
            continue

        # Skip table headers (contain multiple |)
        if text.count('|') >= 3:
            continue

        # Check for section headers like "2.\tScope & Planning"
        section_match = re.match(r'^\d+\.\s+(.+)$', text)
        if section_match:
            title = clean_section_title(text)
            return title, content_list

        # First meaningful text item is the title
        if len(text) > 1 and len(text) < 120:
            remaining = content_list.copy()
            remaining.pop(i)
            return text, remaining

    return "", content_list


def extract_pptx(filepath):
    """Extract all slide content from a PPTX file."""
    prs = Presentation(filepath)
    slides_data = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_info = {
            "slide_number": slide_num,
            "title": "",
            "content": [],
            "notes": ""
        }

        # Extract title from title placeholder
        if slide.shapes.title:
            slide_info["title"] = slide.shapes.title.text.strip()

        # Extract all text from shapes
        all_texts = []
        for shape in slide.shapes:
            texts = extract_text_from_shape(shape)
            all_texts.extend(texts)

        # Filter out standalone page numbers
        filtered_texts = [t for t in all_texts if not is_page_number(t)]

        # Remove vertical tab characters
        filtered_texts = [t.replace('\x0b', ' ').strip() for t in filtered_texts]

        # Remove empty strings
        filtered_texts = [t for t in filtered_texts if t]

        # If no title placeholder, detect from content
        if not slide_info["title"] and filtered_texts:
            detected_title, remaining_content = detect_title(filtered_texts)
            slide_info["title"] = detected_title
            slide_info["content"] = remaining_content
        else:
            slide_info["content"] = filtered_texts

        # Remove title from content if it appears there too
        if slide_info["title"] and slide_info["title"] in slide_info["content"]:
            slide_info["content"].remove(slide_info["title"])

        # Also clean section headers from content (remove numbering)
        cleaned_content = []
        for item in slide_info["content"]:
            # Clean section headers in content
            section_match = re.match(r'^\d+\.\s+(.+)$', item)
            if section_match:
                cleaned = clean_section_title(item)
                if cleaned and cleaned != slide_info["title"]:
                    cleaned_content.append(cleaned)
            else:
                cleaned_content.append(item)
        slide_info["content"] = cleaned_content

        # Extract speaker notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_info["notes"] = notes_text

        slides_data.append(slide_info)

    return slides_data


def build_qa_pairs(slides_data):
    """Build Q&A pairs from slide content for the chatbot.
    Group continuation slides under the same topic.
    """
    # First pass: group slides by topic
    topic_groups = {}
    topic_order = []  # Preserve insertion order
    current_topic = ""

    for slide in slides_data:
        title = slide["title"]
        content = slide["content"]

        if not title and not content:
            continue

        # Determine the topic - always clean "(Cont'd)" to merge with parent
        if title:
            clean_title = clean_section_title(title)
            topic_key = clean_title if clean_title else title
            current_topic = topic_key
        else:
            topic_key = current_topic if current_topic else f"Slide {slide['slide_number']}"

        if topic_key not in topic_groups:
            topic_groups[topic_key] = {
                "first_slide": slide["slide_number"],
                "slides": [],
                "all_content": [],
                "notes": []
            }
            topic_order.append(topic_key)

        topic_groups[topic_key]["slides"].append(slide["slide_number"])
        topic_groups[topic_key]["all_content"].extend(content)
        if slide.get("notes"):
            topic_groups[topic_key]["notes"].append(slide["notes"])

    # Second pass: build QA pairs in original order
    qa_pairs = []
    for topic in topic_order:
        data = topic_groups[topic]
        # Remove duplicate content lines
        seen = set()
        unique_content = []
        for line in data["all_content"]:
            if line not in seen:
                seen.add(line)
                unique_content.append(line)

        answer_parts = []
        if unique_content:
            answer_parts.append("\n".join(unique_content))
        if data["notes"]:
            answer_parts.append("\nAdditional Notes: " + " ".join(data["notes"]))

        full_answer = "\n".join(answer_parts) if answer_parts else "No detailed content available for this topic."

        # Build keywords
        keywords = []
        clean_topic = re.sub(r'[^\w\s]', ' ', topic.lower())
        keywords.extend([w for w in clean_topic.split() if len(w) > 2])

        # Extract top keywords from content
        if unique_content:
            all_text = " ".join(unique_content).lower()
            clean_text = re.sub(r'[^\w\s]', ' ', all_text)
            content_words = [w for w in clean_text.split() if len(w) > 3]
            word_freq = {}
            for w in content_words:
                word_freq[w] = word_freq.get(w, 0) + 1
            top_words = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)[:15]
            keywords.extend([w for w, _ in top_words])

        keywords = list(dict.fromkeys(keywords))

        slide_range = f"{data['first_slide']}"
        if len(data["slides"]) > 1:
            slide_range = f"{data['slides'][0]}-{data['slides'][-1]}"

        qa_pairs.append({
            "slide": slide_range,
            "topic": topic,
            "keywords": keywords,
            "answer": full_answer
        })

    return qa_pairs


def generate_js_file(slides_data, qa_pairs, output_path="npd_data.js"):
    """Generate the JavaScript data file for the chatbot."""

    # Build topic list for menu (unique topics only)
    topics = []
    seen_topics = set()
    for qa in qa_pairs:
        if qa["topic"] not in seen_topics and qa["topic"]:
            topics.append({
                "slide": qa["slide"],
                "title": qa["topic"]
            })
            seen_topics.add(qa["topic"])

    js_content = f"""// Auto-generated from NPD.pptx - DO NOT EDIT MANUALLY
// Generated by extract_pptx.py

const NPD_TOPICS = {json.dumps(topics, indent=2, ensure_ascii=False)};

const NPD_SLIDES = {json.dumps(slides_data, indent=2, ensure_ascii=False)};

const NPD_QA = {json.dumps(qa_pairs, indent=2, ensure_ascii=False)};
"""

    with open(output_path, "w", encoding="utf-8") as f:
        f.write(js_content)

    print(f"  Extracted {len(slides_data)} slides")
    print(f"  Generated {len(qa_pairs)} unique topic Q&A entries")
    print(f"  Found {len(topics)} browsable topics")
    print(f"  Saved to {output_path}")


if __name__ == "__main__":
    import os

    script_dir = os.path.dirname(os.path.abspath(__file__))
    pptx_path = os.path.join(script_dir, "NPD.pptx")
    output_path = os.path.join(script_dir, "npd_data.js")

    if not os.path.exists(pptx_path):
        print(f"ERROR: NPD.pptx not found at {pptx_path}")
        exit(1)

    print("Extracting content from NPD.pptx...")
    slides_data = extract_pptx(pptx_path)
    qa_pairs = build_qa_pairs(slides_data)
    generate_js_file(slides_data, qa_pairs, output_path)
    print("Done! Open index.html in a browser to use the chatbot.")
